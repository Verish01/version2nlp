from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django import forms
from .models import Inventory
import csv
import os
import sqlite3
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
from pptx import Presentation
from pptx.util import Inches
from docx import Document

# Define InventoryForm class
class InventoryForm(forms.ModelForm):
    class Meta:
        model = Inventory
        fields = ['product_name', 'quantity_in_stock', 'cost_per_item', 'quantity_sold', 'sales', 'stock_date', 'photos']
        widgets = {
            'stock_date': forms.DateInput(attrs={'type': 'date'}),
        }

# Define QueryForm class
class QueryForm(forms.Form):
    query = forms.CharField(label='Enter your natural language query', max_length=200)

# Load the ML model and tokenizer
model_path = 'gaussalgo/T5-LM-Large-text2sql-spider'
model = AutoModelForSeq2SeqLM.from_pretrained(model_path)
tokenizer = AutoTokenizer.from_pretrained(model_path)

def query_database(db_path, query):
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        cursor.execute(query)
        results = cursor.fetchall()
        conn.close()
        return results
    except sqlite3.Error as e:
        return str(e)

def get_sql_query(question, schema):
    input_text = f"Question: {question} Schema: {schema}"
    model_inputs = tokenizer(input_text, return_tensors="pt")
    outputs = model.generate(**model_inputs, max_length=512)
    output_text = tokenizer.batch_decode(outputs, skip_special_tokens=True)
    return output_text[0]

def download_csv(results):
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="results.csv"'
    writer = csv.writer(response)
    default_headers = ['Product ID', 'Product Name', 'Quantity in Stock', 'Cost per Item']
    writer.writerow(default_headers)
    for row in results or []:
        writer.writerow(row)
    return response

def download_word(results):
    document = Document()
    document.add_heading('Query Results', 0)
    column_names = [str(col) for col in results[0]] if results else []
    num_cols = len(column_names)
    table = document.add_table(rows=1, cols=num_cols)
    hdr_cells = table.rows[0].cells
    for i, column_name in enumerate(column_names):
        hdr_cells[i].text = column_name
    for row in results[1:]:
        row_cells = table.add_row().cells
        for i, value in enumerate(row):
            row_cells[i].text = str(value)
    document_path = 'query_results.docx'
    document.save(document_path)
    with open(document_path, 'rb') as docx_file:
        response = HttpResponse(docx_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = 'attachment; filename=query_results.docx'
    return response

def download_ppt(results):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Query Results"
    if not results:
        return HttpResponse("No results found")
    num_cols = len(results[0]) if results else 0
    columns = ['Product ID', 'Product Name', 'Quantity in Stock', 'Cost per Item']
    columns = columns[:num_cols]
    rows = len(results) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8)
    table = slide.shapes.add_table(rows, num_cols, left, top, width, height).table
    for col in range(min(num_cols, len(columns))):
        table.cell(0, col).text = columns[col]
    for i, row in enumerate(results or []):
        for j, value in enumerate(row or []):
            table.cell(i + 1, j).text = str(value)
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="results.pptx"'
    prs.save(response)
    return response
def clear_history(request):
    if 'query_history' in request.session:
        del request.session['query_history']
    return redirect('query_view')


def query_view(request):
    results = None
    sql_query = None
    query_history = request.session.get('query_history', [])
    if request.method == 'POST':
        form = QueryForm(request.POST)
        if form.is_valid():
            question = form.cleaned_data['query']
            schema = '''"myapps_inventory" "DATABASE" 
                        "product_ID" int, 
                        "product_name" text, 
                        "quantity_in_stock" int, 
                        "cost_per_item" float, 
                        primary key: "product_ID"'''
            sql_query = get_sql_query(question, schema)
            db_path = os.path.join('ims', 'db.sqlite3')
            results = query_database(db_path, sql_query)
            query_history.append({'question': question, 'sql_query': sql_query, 'results': results})
            request.session['query_history'] = query_history
            if 'csv' in request.POST:
                return download_csv(results)
            elif 'word' in request.POST:
                return download_word(results)
            elif 'ppt' in request.POST:
                return download_ppt(results)
            response_data = {'form': form, 'results': results, 'query': sql_query, 'query_history': query_history}
            return render(request, 'query_view.html', response_data)
    else:
        form = QueryForm()
    return render(request, 'query_view.html', {'form': form, 'results': results, 'query': sql_query, 'query_history': query_history})

def inventory_list(request):
    inventory_items = Inventory.objects.all()
    return render(request, 'inventory_list.html', {'inventory_items': inventory_items})

def edit_item(request, id):
    item = get_object_or_404(Inventory, id=id)
    if request.method == "POST":
        form = InventoryForm(request.POST, instance=item)
        if form.is_valid():
            form.save()
            return redirect('inventory_list')
    else:
        form = InventoryForm(instance=item)
    return render(request, 'edit_item.html', {'form': form})

def delete_item(request, id):
    item = get_object_or_404(Inventory, id=id)
    if request.method == "POST":
        item.delete()
        return redirect('inventory_list')
    return render(request, 'confirm_delete.html', {'item': item})

def add_item(request):
    if request.method == 'POST':
        form = InventoryForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect('inventory_list')
    else:
        form = InventoryForm()
    return render(request, 'add_item.html', {'form': form})

def inventory_visualization(request):
    return render(request, 'gradio.html')

def export_inventory_csv(request):
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="inventory.csv"'
    writer = csv.writer(response)
    writer.writerow(['Product Name', 'Quantity in Stock', 'Cost per Item', 'Quantity Sold', 'Sales', 'Stock Date'])
    for item in Inventory.objects.all():
        writer.writerow([item.product_name, item.quantity_in_stock, item.cost_per_item, item.quantity_sold, item.sales, item.stock_date])
    return response
